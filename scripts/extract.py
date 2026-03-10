#!/usr/bin/env python3
"""
Extract slides and metadata from a PPTX presentation.

Exports each slide as a JPG image and extracts text content, speaker notes,
and layout metadata into a structured JSON file.

Usage:
    python3 extract.py presentation.pptx [--output-dir .]

Requirements:
    - python-pptx: pip install python-pptx
    - LibreOffice: brew install --cask libreoffice (macOS) or apt install libreoffice (Linux)

Output:
    slides/slide-001.jpg, slide-002.jpg, ...   (one per slide)
    slides-data.json                            (structured metadata)
"""

import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile


def check_dependencies():
    """Verify required dependencies are available."""
    errors = []

    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        errors.append(
            "python-pptx is not installed.\n"
            "  Install with: pip install python-pptx"
        )

    if not shutil.which("soffice") and not shutil.which("libreoffice"):
        errors.append(
            "LibreOffice is not installed (needed to render slides as images).\n"
            "  macOS:  brew install --cask libreoffice\n"
            "  Linux:  sudo apt install libreoffice"
        )

    if errors:
        print("Missing dependencies:\n", file=sys.stderr)
        for e in errors:
            print(f"  {e}\n", file=sys.stderr)
        sys.exit(1)


def export_slide_images(pptx_path, slides_dir):
    """Use LibreOffice to export each slide as a JPG."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")

    with tempfile.TemporaryDirectory() as tmpdir:
        # LibreOffice converts PPTX to a set of JPGs
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to", "jpg",
                "--outdir", tmpdir,
                os.path.abspath(pptx_path),
            ],
            capture_output=True,
            text=True,
        )

        if result.returncode != 0:
            print(f"LibreOffice export failed:\n{result.stderr}", file=sys.stderr)
            sys.exit(1)

        # LibreOffice exports a single JPG for single-image conversion.
        # For multi-slide, we need PDF intermediary then convert pages.
        # Check what we got:
        exported = sorted(f for f in os.listdir(tmpdir) if f.endswith(".jpg"))

        if len(exported) == 1:
            # LibreOffice exported a single image — need PDF route instead
            return _export_via_pdf(pptx_path, slides_dir, soffice, tmpdir)

        # Multiple JPGs — rename to slide-NNN.jpg format
        os.makedirs(slides_dir, exist_ok=True)
        for i, fname in enumerate(exported, 1):
            src = os.path.join(tmpdir, fname)
            dst = os.path.join(slides_dir, f"slide-{i:03d}.jpg")
            shutil.copy2(src, dst)

        return len(exported)


def _export_via_pdf(pptx_path, slides_dir, soffice, tmpdir):
    """Export slides via PDF intermediary (more reliable for multi-slide decks)."""
    # Step 1: Convert PPTX to PDF
    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", tmpdir,
            os.path.abspath(pptx_path),
        ],
        capture_output=True,
        text=True,
    )

    if result.returncode != 0:
        print(f"LibreOffice PDF export failed:\n{result.stderr}", file=sys.stderr)
        sys.exit(1)

    pdf_files = [f for f in os.listdir(tmpdir) if f.endswith(".pdf")]
    if not pdf_files:
        print("Error: LibreOffice produced no PDF output.", file=sys.stderr)
        sys.exit(1)

    pdf_path = os.path.join(tmpdir, pdf_files[0])

    # Step 2: Convert PDF pages to JPGs using sips (macOS) or ImageMagick
    return _pdf_to_jpgs(pdf_path, slides_dir)


def _pdf_to_jpgs(pdf_path, slides_dir):
    """Convert PDF pages to individual JPG files."""
    os.makedirs(slides_dir, exist_ok=True)

    # Try pdftoppm first (poppler-utils, cross-platform)
    if shutil.which("pdftoppm"):
        with tempfile.TemporaryDirectory() as img_tmp:
            prefix = os.path.join(img_tmp, "slide")
            subprocess.run(
                ["pdftoppm", "-jpeg", "-r", "150", pdf_path, prefix],
                capture_output=True,
                text=True,
                check=True,
            )
            exported = sorted(f for f in os.listdir(img_tmp) if f.endswith(".jpg"))
            for i, fname in enumerate(exported, 1):
                shutil.copy2(
                    os.path.join(img_tmp, fname),
                    os.path.join(slides_dir, f"slide-{i:03d}.jpg"),
                )
            return len(exported)

    # Fallback: Python-based conversion
    try:
        from pdf2image import convert_from_path

        images = convert_from_path(pdf_path, dpi=150, fmt="jpeg")
        for i, img in enumerate(images, 1):
            img.save(os.path.join(slides_dir, f"slide-{i:03d}.jpg"), "JPEG")
        return len(images)
    except ImportError:
        pass

    # Fallback: macOS sips (can convert PDF but only first page)
    # Not viable for multi-page — error out with guidance
    print(
        "Error: Cannot convert PDF pages to images.\n"
        "Install one of:\n"
        "  poppler (recommended): brew install poppler   (provides pdftoppm)\n"
        "  pdf2image:             pip install pdf2image",
        file=sys.stderr,
    )
    sys.exit(1)


def extract_metadata(pptx_path):
    """Extract text content, speaker notes, and layout info from each slide."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(pptx_path)

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        # Extract text from all shapes
        title = ""
        text_content = []
        has_images = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # First text shape with content is likely the title
                if not title and (shape.shape_type == 13 or hasattr(shape, "is_placeholder")):
                    try:
                        pf = shape.placeholder_format
                        if pf is not None and pf.idx == 0:
                            title = text
                            continue
                    except (ValueError, AttributeError):
                        pass
                text_content.append(text)

            if shape.shape_type == 13:  # Picture
                has_images = True

        # If no title found via placeholder, use first text
        if not title and text_content:
            title = text_content.pop(0)

        # Speaker notes
        speaker_notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            speaker_notes = slide.notes_slide.notes_text_frame.text.strip()

        # Layout name
        layout_name = slide.slide_layout.name if slide.slide_layout else ""

        slides.append({
            "number": i,
            "image": f"slides/slide-{i:03d}.jpg",
            "title": title,
            "text_content": text_content,
            "speaker_notes": speaker_notes,
            "has_images": has_images,
            "layout_name": layout_name,
        })

    # Presentation metadata
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    has_any_notes = any(s["speaker_notes"] for s in slides)

    return {
        "slides": slides,
        "metadata": {
            "slide_count": len(slides),
            "dimensions": {
                "width": slide_width,
                "height": slide_height,
            },
            "has_speaker_notes": has_any_notes,
        },
    }


def main():
    parser = argparse.ArgumentParser(
        description="Extract slides and metadata from a PPTX presentation."
    )
    parser.add_argument("pptx", help="Path to .pptx file")
    parser.add_argument(
        "--output-dir", "-o", default=".",
        help="Output directory (default: current directory)"
    )
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print(f"Error: File not found: {args.pptx}", file=sys.stderr)
        sys.exit(1)

    if not args.pptx.endswith(".pptx"):
        print(f"Error: Expected a .pptx file, got: {args.pptx}", file=sys.stderr)
        sys.exit(1)

    check_dependencies()

    output_dir = os.path.abspath(args.output_dir)
    slides_dir = os.path.join(output_dir, "slides")
    json_path = os.path.join(output_dir, "slides-data.json")

    # Export slide images
    print(f"Exporting slide images to {slides_dir}/...", file=sys.stderr)
    count = export_slide_images(args.pptx, slides_dir)
    print(f"Exported {count} slide images.", file=sys.stderr)

    # Extract metadata
    print("Extracting text and metadata...", file=sys.stderr)
    data = extract_metadata(args.pptx)

    # Verify slide count matches
    if count != data["metadata"]["slide_count"]:
        print(
            f"Warning: Image export produced {count} images but PPTX has "
            f"{data['metadata']['slide_count']} slides.",
            file=sys.stderr,
        )

    # Write JSON
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

    print(f"Wrote {json_path} ({data['metadata']['slide_count']} slides)", file=sys.stderr)


if __name__ == "__main__":
    main()
